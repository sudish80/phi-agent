"""File handler module for J.A.R.V.I.S.

Handles reading, writing, and converting between file formats:
  - CSV, Excel (XLSX), PDF, JSON, YAML, XML, DOCX, PPTX, Markdown
  - SVG image generation
  - LaTeX compilation to PDF
"""

import asyncio
import csv
import io
import json
import logging
import os
import re
import subprocess
import tempfile
import uuid
from pathlib import Path
from typing import Optional, Dict, Any, List

from backend.shared.config import settings

logger = logging.getLogger(__name__)

OUTPUT_DIR = Path(__file__).resolve().parent / "static" / "files"
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
_BASE_URL = f"http://localhost:{settings.action_port}/static/files"


# ─── CSV ───────────────────────────────────────────────────────────────────

async def read_csv(file_path: str, delimiter: str = ",", has_header: bool = True,
                   max_rows: int = 100) -> str:
    """Read a CSV file and return formatted contents."""
    loop = asyncio.get_event_loop()

    def _read():
        full = os.path.abspath(os.path.expanduser(file_path))
        if not os.path.exists(full):
            return f"File not found: {full}"
        rows = []
        with open(full, "r", encoding="utf-8", errors="replace") as f:
            reader = csv.reader(f, delimiter=delimiter)
            for i, row in enumerate(reader):
                if i >= max_rows + (1 if has_header else 0):
                    rows.append(f"... [{i - (1 if has_header else 0)} more rows]")
                    break
                rows.append(f"  Row {i}: {' | '.join(row)}")
        header = rows.pop(0) if has_header and rows else ""
        out = f"**{os.path.basename(full)}** ({len(rows)} rows shown)"
        if header:
            out += f"\n  Header: {header}"
        out += "\n" + "\n".join(rows)
        return out

    return await loop.run_in_executor(None, _read)


async def query_csv(file_path: str, column: str, value: str,
                    delimiter: str = ",") -> str:
    """Filter CSV rows where column equals value."""
    loop = asyncio.get_event_loop()

    def _query():
        full = os.path.abspath(os.path.expanduser(file_path))
        if not os.path.exists(full):
            return f"File not found: {full}"
        with open(full, "r", encoding="utf-8", errors="replace") as f:
            reader = csv.DictReader(f, delimiter=delimiter)
            if column not in reader.fieldnames:
                return f"Column '{column}' not found. Columns: {', '.join(reader.fieldnames)}"
            matches = []
            for row in reader:
                if row.get(column) == value:
                    matches.append(" | ".join(f"{k}={v}" for k, v in row.items()))
            if not matches:
                return f"No rows where {column} = '{value}'"
            return f"**{len(matches)} match(es)**\n" + "\n".join(f"  {m}" for m in matches[:50])

    return await loop.run_in_executor(None, _query)


# ─── Excel (XLSX) ───────────────────────────────────────────────────────────

async def read_excel(file_path: str, sheet: str = None) -> str:
    """Read an Excel file and return formatted contents."""
    try:
        import pandas as pd
    except ImportError:
        return "Excel support requires pandas + openpyxl"

    loop = asyncio.get_event_loop()

    def _read():
        full = os.path.abspath(os.path.expanduser(file_path))
        if not os.path.exists(full):
            return f"File not found: {full}"
        xl = pd.ExcelFile(full)
        sheets = xl.sheet_names
        if sheet and sheet not in sheets:
            return f"Sheet '{sheet}' not found. Available: {', '.join(sheets)}"
        target = sheet or sheets[0]
        df = pd.read_excel(full, sheet_name=target)
        out = f"**{os.path.basename(full)}** → Sheet: **{target}**"
        out += f"\n{df.shape[0]} rows × {df.shape[1]} cols\n"
        out += df.head(50).to_string(index=False)
        if df.shape[0] > 50:
            out += f"\n... [{df.shape[0] - 50} more rows]"
        return out

    return await loop.run_in_executor(None, _read)


async def write_excel(data: Dict[str, List[Dict]], filename: str = None) -> str:
    """Write data to an Excel file. data = {sheet_name: [{col: val}, ...]}."""
    try:
        import pandas as pd
    except ImportError:
        return "Excel write requires pandas + openpyxl"

    loop = asyncio.get_event_loop()

    def _write():
        name = filename or f"export_{uuid.uuid4().hex[:8]}.xlsx"
        path = OUTPUT_DIR / name
        with pd.ExcelWriter(path, engine="openpyxl") as writer:
            for sheet_name, records in data.items():
                df = pd.DataFrame(records)
                df.to_excel(writer, sheet_name=sheet_name, index=False)
        return f"{_BASE_URL}/{name}"

    return await loop.run_in_executor(None, _write)


# ─── PDF ────────────────────────────────────────────────────────────────────

async def read_pdf(file_path: str, max_pages: int = 20) -> str:
    """Extract text from a PDF file."""
    loop = asyncio.get_event_loop()

    def _read():
        full = os.path.abspath(os.path.expanduser(file_path))
        if not os.path.exists(full):
            return f"File not found: {full}"
        try:
            import fitz
            doc = fitz.open(full)
            out = f"**{os.path.basename(full)}** — {doc.page_count} pages\n"
            for i, page in enumerate(doc):
                if i >= max_pages:
                    out += f"\n... [{doc.page_count - max_pages} more pages]"
                    break
                text = page.get_text().strip()
                if text:
                    out += f"\n--- Page {i + 1} ---\n{text[:3000]}"
                else:
                    out += f"\n--- Page {i + 1} --- (no text)"
                if i >= max_pages - 1:
                    break
            doc.close()
            return out
        except ImportError:
            return "PDF support requires PyMuPDF (fitz)"
        except Exception as e:
            return f"PDF error: {e}"

    return await loop.run_in_executor(None, _read)


async def pdf_metadata(file_path: str) -> str:
    """Get PDF metadata."""
    loop = asyncio.get_event_loop()

    def _meta():
        full = os.path.abspath(os.path.expanduser(file_path))
        import fitz
        doc = fitz.open(full)
        meta = doc.metadata
        doc.close()
        return "\n".join(f"  {k}: {v}" for k, v in meta.items() if v)

    try:
        return await loop.run_in_executor(None, _meta)
    except ImportError:
        return "PDF support requires PyMuPDF"
    except Exception as e:
        return f"PDF metadata error: {e}"


# ─── JSON ───────────────────────────────────────────────────────────────────

async def read_json(file_path: str, pretty: bool = True) -> str:
    """Read and format a JSON file."""
    loop = asyncio.get_event_loop()

    def _read():
        full = os.path.abspath(os.path.expanduser(file_path))
        with open(full, "r", encoding="utf-8") as f:
            data = json.load(f)
        return json.dumps(data, indent=2) if pretty else str(data)

    try:
        result = await loop.run_in_executor(None, _read)
        return result[:15000] + ("\n... [truncated]" if len(result) > 15000 else "")
    except Exception as e:
        return f"JSON error: {e}"


async def query_json(file_path: str, key_path: str) -> str:
    """Query a JSON file using dot-separated keys (e.g. 'data.users.0.name')."""
    loop = asyncio.get_event_loop()

    def _query():
        full = os.path.abspath(os.path.expanduser(file_path))
        with open(full, "r", encoding="utf-8") as f:
            data = json.load(f)
        parts = key_path.split(".")
        current = data
        for p in parts:
            if isinstance(current, dict):
                current = current.get(p)
            elif isinstance(current, list):
                try:
                    current = current[int(p)]
                except (IndexError, ValueError):
                    return f"Key '{p}' not found in list (index {p})"
            else:
                return f"Cannot traverse: {type(current)} at '{p}'"
            if current is None:
                return f"Key '{p}' not found"
        return json.dumps(current, indent=2)[:5000]

    return await loop.run_in_executor(None, _query)


# ─── YAML ───────────────────────────────────────────────────────────────────

async def read_yaml(file_path: str) -> str:
    """Read and format a YAML file."""
    try:
        import yaml
    except ImportError:
        return "YAML support requires pyyaml"

    loop = asyncio.get_event_loop()

    def _read():
        full = os.path.abspath(os.path.expanduser(file_path))
        with open(full, "r", encoding="utf-8") as f:
            data = yaml.safe_load(f)
        return yaml.dump(data, default_flow_style=False, allow_unicode=True)[:10000]

    return await loop.run_in_executor(None, _read)


# ─── XML ────────────────────────────────────────────────────────────────────

async def read_xml(file_path: str) -> str:
    """Read and pretty-print an XML file."""
    loop = asyncio.get_event_loop()

    def _read():
        full = os.path.abspath(os.path.expanduser(file_path))
        from xml.dom import minidom
        with open(full, "r", encoding="utf-8") as f:
            dom = minidom.parse(f)
        return dom.toprettyxml(indent="  ")[:10000]

    return await loop.run_in_executor(None, _read)


# ─── DOCX (Word) ────────────────────────────────────────────────────────────

async def read_docx(file_path: str) -> str:
    """Extract text from a Word document."""
    try:
        import docx
    except ImportError:
        return "DOCX support requires python-docx"

    loop = asyncio.get_event_loop()

    def _read():
        full = os.path.abspath(os.path.expanduser(file_path))
        doc = docx.Document(full)
        paras = [p.text for p in doc.paragraphs if p.text.strip()]
        return f"**{os.path.basename(full)}** — {len(paras)} paragraphs\n" + "\n".join(paras[:200])

    return await loop.run_in_executor(None, _read)


# ─── Markdown ───────────────────────────────────────────────────────────────

async def read_markdown(file_path: str) -> str:
    """Read a Markdown file and return structured outline."""
    loop = asyncio.get_event_loop()

    def _read():
        full = os.path.abspath(os.path.expanduser(file_path))
        with open(full, "r", encoding="utf-8") as f:
            lines = f.readlines()
        toc = []
        for line in lines:
            m = re.match(r"^(#{1,6})\s+(.+)", line)
            if m:
                level = len(m.group(1))
                toc.append(f"{'  ' * (level - 1)}* {m.group(2).strip()}")
        content = "".join(lines)
        return f"**{os.path.basename(full)}** — {len(lines)} lines\n\n## Outline\n" + "\n".join(
            toc[:50]) + "\n\n## Content (first 5000 chars)\n" + content[:5000]

    return await loop.run_in_executor(None, _read)


# ─── SVG Generation ─────────────────────────────────────────────────────────

async def generate_svg(shape: str = "circle", width: int = 400, height: int = 300,
                       color: str = "#007AFF", label: str = "",
                       svg_content: str = None, filename: str = None) -> str:
    """Generate an SVG image. Can use preset shapes or raw SVG content."""
    name = (filename if filename and filename.endswith(".svg")
             else f"{filename or 'svg_' + uuid.uuid4().hex[:8]}.svg")
    path = OUTPUT_DIR / name

    if svg_content:
        with open(path, "w", encoding="utf-8") as f:
            f.write(svg_content)
        return f"{_BASE_URL}/{name}"

    try:
        import svgwrite
    except ImportError:
        fallback = (
            f'<svg xmlns="http://www.w3.org/2000/svg" width="{width}" height="{height}">'
            f'<rect width="100%" height="100%" fill="#f8f9fa"/>'
        )
        if shape == "circle":
            fallback += f'<circle cx="{width // 2}" cy="{height // 2}" r="{min(width, height) // 3}" fill="{color}"/>'
        elif shape == "square":
            s = min(width, height) // 2
            fallback += f'<rect x="{(width - s) // 2}" y="{(height - s) // 2}" width="{s}" height="{s}" fill="{color}"/>'
        elif shape == "triangle":
            cx, cy = width // 2, height // 3
            r = min(width, height) // 3
            pts = f"{cx},{cy - r} {cx - r},{cy + r} {cx + r},{cy + r}"
            fallback += f'<polygon points="{pts}" fill="{color}"/>'
        if label:
            fallback += f'<text x="{width // 2}" y="{height - 20}" text-anchor="middle" font-family="Arial" font-size="16" fill="#333">{label}</text>'
        fallback += "</svg>"
        with open(path, "w", encoding="utf-8") as f:
            f.write(fallback)
        return f"{_BASE_URL}/{name}"

    dwg = svgwrite.Drawing(str(path), size=(width, height))
    dwg.add(dwg.rect(insert=(0, 0), size=("100%", "100%"), fill="#f8f9fa"))

    cx, cy = width // 2, height // 2
    r = min(width, height) // 3
    if shape == "circle":
        dwg.add(dwg.circle(center=(cx, cy), r=r, fill=color))
    elif shape == "square":
        dwg.add(dwg.rect(insert=(cx - r, cy - r), size=(r * 2, r * 2), fill=color))
    elif shape == "triangle":
        dwg.add(dwg.polygon(points=[(cx, cy - r), (cx - r, cy + r), (cx + r, cy + r)], fill=color))
    elif shape == "line_chart":
        pts = [(50 + i * 60, 250 - v) for i, v in enumerate([50, 120, 80, 190, 140, 200, 180])]
        dwg.add(dwg.polyline(points=pts, stroke=color, fill="none", stroke_width=3))
        for px, py in pts:
            dwg.add(dwg.circle(center=(px, py), r=4, fill=color))
    elif shape == "bar_chart":
        bw = (width - 100) // 5
        for i, v in enumerate([60, 120, 90, 150, 110]):
            x = 50 + i * (bw + 10)
            dwg.add(dwg.rect(insert=(x, 250 - v), size=(bw, v), fill=color, opacity=0.8))
    elif shape == "pie_chart":
        vals = [30, 25, 20, 15, 10]
        colors = ["#FF6B6B", "#4ECDC4", "#45B7D1", "#FFEAA7", "#DDA0DD"]
        start = 0
        for v, c in zip(vals, colors):
            angle = 360 * v / sum(vals)
            dwg.add(dwg.path(d=f"M {cx},{cy} L {cx + r},{cy} A {r},{r} 0 0 1 {cx + r * 0.5},{cy - r * 0.866} Z",
                             fill=c))
    if label:
        dwg.add(dwg.text(label, insert=(cx, height - 15), text_anchor="middle",
                         font_family="Arial", font_size="14", fill="#333"))
    dwg.save()
    return f"{_BASE_URL}/{name}"


# ─── LaTeX Compilation ──────────────────────────────────────────────────────

async def write_latex(content: str, filename: str = None) -> str:
    """Write a .tex file and optionally compile to PDF.

    Returns the URL to the .tex file (and .pdf if compilation succeeds).
    """
    name = filename or f"latex_{uuid.uuid4().hex[:8]}"
    tex_path = OUTPUT_DIR / f"{name}.tex"

    with open(tex_path, "w", encoding="utf-8") as f:
        f.write(content)

    tex_url = f"{_BASE_URL}/{name}.tex"
    return f"LaTeX source saved: {tex_url}\n\nTo compile, use compile_latex with this filename."


async def compile_latex(filename: str) -> str:
    """Compile a .tex file to PDF using pdflatex.

    The filename should be the name (without .tex) saved by write_latex.
    """
    name = filename.replace(".tex", "")
    tex_path = OUTPUT_DIR / f"{name}.tex"
    pdf_path = OUTPUT_DIR / f"{name}.pdf"

    if not tex_path.exists():
        # Search for any matching file
        matches = list(OUTPUT_DIR.glob(f"{name}*"))
        if matches:
            tex_path = matches[0]
        else:
            return f"File '{name}.tex' not found. Use write_latex first."

    try:
        result = subprocess.run(
            ["pdflatex", "-interaction=nonstopmode", "-output-directory",
             str(OUTPUT_DIR), str(tex_path)],
            capture_output=True, text=True, timeout=60,
        )
        if pdf_path.exists():
            return f"PDF compiled: {_BASE_URL}/{name}.pdf"
        else:
            log = result.stdout[-1000:] if result.stdout else result.stderr[-1000:]
            return f"LaTeX compilation failed. Log:\n{log}"
    except FileNotFoundError:
        return "pdflatex not found. Install TeX distribution (MiKTeX / TeX Live) to compile."
    except subprocess.TimeoutExpired:
        return "LaTeX compilation timed out (60s)."
    except Exception as e:
        return f"LaTeX compilation error: {e}"


# ─── PowerPoint (PPTX) ─────────────────────────────────────────────────────

async def create_pptx(slides: List[Dict], title: str = "Presentation",
                      filename: str = None) -> str:
    """Create a PowerPoint presentation from slide definitions.

    Each slide dict can have:
      - title (str): slide title
      - content (str | List[str]): body text or bullet points
      - subtitle (str): optional subtitle
      - layout (str): 'title', 'content', 'two_content', 'blank', 'title_only'
      - table (Dict): {headers: [...], rows: [[...], ...]}
      - chart (Dict): {type: 'bar'|'line'|'pie', title: str, labels: [...], values: [...], values2: [...]}
      - notes (str): speaker notes
      - image (str): URL or local path to image
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE
    except ImportError:
        return "PPTX support requires python-pptx"

    loop = asyncio.get_event_loop()

    def _build():
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        layout_map = {
            "title": prs.slide_layouts[0],
            "content": prs.slide_layouts[1],
            "two_content": prs.slide_layouts[3],
            "blank": prs.slide_layouts[6],
            "title_only": prs.slide_layouts[5],
        }

        for slide_def in slides:
            layout = layout_map.get(slide_def.get("layout", "content"), prs.slide_layouts[1])
            slide = prs.slides.add_slide(layout)

            slide_title = slide_def.get("title", "")
            if slide.shapes.title:
                slide.shapes.title.text = slide_title

            content = slide_def.get("content", "")
            subtitle = slide_def.get("subtitle", "")

            if subtitle and slide.placeholders:
                for ph in slide.placeholders:
                    if ph.placeholder_format.idx == 1:
                        ph.text = subtitle
                        break

            if content and slide.placeholders:
                for ph in slide.placeholders:
                    idx = ph.placeholder_format.idx
                    if isinstance(content, list):
                        if idx in (1, 2):
                            tf = ph.text_frame
                            tf.clear()
                            for i, item in enumerate(content):
                                if i == 0:
                                    tf.text = str(item)
                                else:
                                    p = tf.add_paragraph()
                                    p.text = str(item)
                                    p.level = 0
                            break
                    else:
                        if idx == 1 and not subtitle:
                            ph.text = str(content)
                            break

            if slide_def.get("chart"):
                c = slide_def["chart"]
                chart_type_map = {
                    "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
                    "line": XL_CHART_TYPE.LINE_MARKERS,
                    "pie": XL_CHART_TYPE.PIE,
                }
                xtype = chart_type_map.get(c.get("type", "bar"))
                chart_data = CategoryChartData()
                chart_data.categories = c.get("labels", [])
                chart_data.add_series(c.get("title", "Series"), c.get("values", []))
                if c.get("values2"):
                    chart_data.add_series("Series 2", c["values2"])
                slide.shapes.add_chart(
                    xtype or XL_CHART_TYPE.COLUMN_CLUSTERED,
                    Inches(1), Inches(2), Inches(6), Inches(4.5), chart_data
                )

            if slide_def.get("table"):
                t = slide_def["table"]
                rows = len(t.get("rows", [])) + 1
                cols = len(t.get("headers", []))
                shape = slide.shapes.add_table(rows, cols, Inches(1), Inches(2.5),
                                                Inches(8), Inches(0.5 * rows))
                table = shape.table
                for i, h in enumerate(t.get("headers", [])):
                    table.cell(0, i).text = str(h)
                for r, row_data in enumerate(t.get("rows", []), 1):
                    for c, val in enumerate(row_data):
                        table.cell(r, c).text = str(val)

            if slide_def.get("notes"):
                notes_slide = slide.notes_slide
                notes_slide.notes_text_frame.text = slide_def["notes"]

            if slide_def.get("image"):
                img_path = slide_def["image"]
                if os.path.exists(img_path):
                    slide.shapes.add_picture(img_path, Inches(1), Inches(2),
                                             Inches(5), Inches(4))

        name = (filename if filename and filename.endswith(".pptx")
                 else f"{filename or 'pres_' + uuid.uuid4().hex[:8]}.pptx")
        path = OUTPUT_DIR / name
        prs.save(str(path))
        return f"{_BASE_URL}/{name}"

    return await loop.run_in_executor(None, _build)


# ─── Obsidian Vault Access ─────────────────────────────────────────────────

def _find_vault(vault_path: str = None) -> Optional[Path]:
    """Find an Obsidian vault directory (contains .obsidian folder)."""
    if vault_path:
        p = Path(os.path.abspath(os.path.expanduser(vault_path)))
        if (p / ".obsidian").is_dir():
            return p
        return None
    # Search common locations
    home = Path.home()
    candidates = [
        home / "Documents" / "Obsidian",
        home / "Documents" / "Obsidian Vault",
        home / "Obsidian",
        home / "Desktop" / "Obsidian",
        home / "Documents" / "My Vault",
        home / "vault",
        Path.cwd(),
    ]
    for c in candidates:
        if (c / ".obsidian").is_dir():
            return c
    return None


async def obsidian_list_vaults() -> str:
    """List all discovered Obsidian vaults on the system."""
    loop = asyncio.get_event_loop()

    def _scan():
        home = Path.home()
        results = []
        for root, dirs, _ in os.walk(str(home), topdown=True):
            if ".obsidian" in dirs:
                results.append(str(Path(root) / ".obsidian"))
                dirs.clear()
            if len(results) >= 10:
                break
            parts = Path(root).relative_to(home).parts if Path(root) != home else []
            if len(parts) > 5:
                dirs.clear()
        if not results:
            return "No Obsidian vaults found."
        lines = ["**Obsidian Vaults Found:**"]
        for r in results:
            vault = Path(r).parent
            lines.append(f"  - {vault.name}: {vault}")
        return "\n".join(lines)

    return await loop.run_in_executor(None, _scan)


async def obsidian_read_note(note_name: str, vault_path: str = None) -> str:
    """Read a note from an Obsidian vault by name (without .md)."""
    vault = _find_vault(vault_path)
    if not vault:
        return "Obsidian vault not found. Use obsidian_list_vaults or provide vault_path."

    loop = asyncio.get_event_loop()

    def _read():
        name = note_name if note_name.endswith(".md") else f"{note_name}.md"
        for f in vault.rglob(name):
            content = f.read_text(encoding="utf-8", errors="replace")
            tags = re.findall(r"#\w+", content)
            links = re.findall(r"\[\[([^\]]+)\]\]", content)
            out = f"**{f.relative_to(vault)}**\n"
            if tags:
                out += f"Tags: {' '.join(tags)}\n"
            if links:
                out += f"Links: {', '.join(links[:20])}\n"
            out += "\n" + content[:8000]
            if len(content) > 8000:
                out += "\n... [truncated]"
            return out
        return f"Note '{note_name}' not found in vault."

    return await loop.run_in_executor(None, _read)


async def obsidian_search_notes(query: str, vault_path: str = None) -> str:
    """Search for notes in an Obsidian vault containing the query text."""
    vault = _find_vault(vault_path)
    if not vault:
        return "Obsidian vault not found."

    loop = asyncio.get_event_loop()

    def _search():
        matches = []
        for f in vault.rglob("*.md"):
            try:
                content = f.read_text(encoding="utf-8", errors="replace")
                if query.lower() in content.lower():
                    preview = ""
                    idx = content.lower().find(query.lower())
                    if idx >= 0:
                        start = max(0, idx - 60)
                        end = min(len(content), idx + len(query) + 60)
                        preview = content[start:end].replace("\n", " ")
                    matches.append(f"  - {f.relative_to(vault)}: ...{preview}..."
                                  if preview else f"  - {f.relative_to(vault)}")
            except Exception:
                continue
        if not matches:
            return f"No notes matching '{query}' found."
        return f"**{len(matches)} notes matching '{query}'**\n" + "\n".join(matches[:30])

    return await loop.run_in_executor(None, _search)


async def obsidian_write_note(note_name: str, content: str,
                              vault_path: str = None) -> str:
    """Write or overwrite a note in an Obsidian vault."""
    vault = _find_vault(vault_path)
    if not vault:
        return "Obsidian vault not found."

    name = note_name if note_name.endswith(".md") else f"{note_name}.md"
    path = vault / name
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(content, encoding="utf-8")
    return f"Note written: {path} ({len(content)} chars)"


async def obsidian_create_note(note_name: str, content: str, tags: str = "",
                               vault_path: str = None) -> str:
    """Create a new note with optional frontmatter tags in an Obsidian vault."""
    if tags:
        tag_list = ", ".join(t.strip().lstrip("#") for t in tags.split(","))
        frontmatter = f"---\ntags: [{tag_list}]\n---\n\n"
        content = frontmatter + content
    return await obsidian_write_note(note_name, content, vault_path)


async def obsidian_list_notes(vault_path: str = None, pattern: str = "*.md") -> str:
    """List all notes in an Obsidian vault."""
    vault = _find_vault(vault_path)
    if not vault:
        return "Obsidian vault not found."

    loop = asyncio.get_event_loop()

    def _list():
        notes = sorted(vault.rglob(pattern))
        if not notes:
            return f"No notes matching {pattern} found."
        lines = [f"**{len(notes)} notes in {vault.name}:**"]
        for n in notes[:100]:
            rel = n.relative_to(vault)
            size = n.stat().st_size
            lines.append(f"  - {rel} ({size} bytes)")
        if len(notes) > 100:
            lines.append(f"  ... and {len(notes) - 100} more")
        return "\n".join(lines)

    return await loop.run_in_executor(None, _list)


# ─── Auto-detect format ─────────────────────────────────────────────────────

async def read_file(file_path: str) -> str:
    """Auto-detect file format by extension and read accordingly."""
    ext = Path(file_path).suffix.lower()
    readers = {
        ".csv": lambda: read_csv(file_path),
        ".tsv": lambda: read_csv(file_path, delimiter="\t"),
        ".xlsx": lambda: read_excel(file_path),
        ".xls": lambda: read_excel(file_path),
        ".pdf": lambda: read_pdf(file_path),
        ".json": lambda: read_json(file_path),
        ".yaml": lambda: read_yaml(file_path),
        ".yml": lambda: read_yaml(file_path),
        ".xml": lambda: read_xml(file_path),
        ".docx": lambda: read_docx(file_path),
        ".md": lambda: read_markdown(file_path),
        ".txt": lambda: _read_text(file_path),
        ".svg": lambda: _read_text(file_path),
        ".tex": lambda: _read_text(file_path),
    }
    reader = readers.get(ext)
    if not reader:
        return f"Unsupported format: {ext}. Supported: {', '.join(readers.keys())}"
    return await reader()


async def _read_text(file_path: str) -> str:
    loop = asyncio.get_event_loop()

    def _read():
        full = os.path.abspath(os.path.expanduser(file_path))
        with open(full, "r", encoding="utf-8", errors="replace") as f:
            return f.read(50000)

    return await loop.run_in_executor(None, _read)
